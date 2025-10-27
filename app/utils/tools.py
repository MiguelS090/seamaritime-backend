# app/utils/tools.py
from io import BytesIO
# import fitz  # PyMuPDF para PDFs - temporariamente comentado
import pandas as pd  # pandas para Excel
import docx  # python-docx para Word
from pptx import Presentation  # python-pptx para PowerPoint

class Tools:
    # Apenas arquivos baseados em texto são permitidos
    ALLOWED_EXTENSIONS = {"pdf", "txt", "xls", "xlsx", "ppt", "pptx", "doc", "docx"}

    def extract_text_from_file(self, file_stream: BytesIO, extension: str) -> str:
        extension = extension.lower()
        if extension not in self.ALLOWED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file extension: {extension}. Only text-based files are allowed."
            )
        if extension == "pdf":
            return self.extract_text_from_pdf(file_stream)
        elif extension in ["xls", "xlsx"]:
            return self.extract_text_from_excel(file_stream)
        elif extension == "txt":
            return self.extract_text_from_txt(file_stream)
        elif extension in ["ppt", "pptx"]:
            return self.extract_text_from_powerpoint(file_stream)
        elif extension in ["doc", "docx"]:
            return self.extract_text_from_doc(file_stream)
        else:
            raise ValueError(f"Unsupported file extension: {extension}")

    def extract_text_from_pdf(self, file_stream: BytesIO) -> str:
        # Temporariamente desabilitado devido a problemas com PyMuPDF
        return "PDF processing temporarily disabled. Please use Azure OCR for PDF processing."
        # text = ""
        # pdf_document = fitz.open(stream=file_stream, filetype="pdf")
        # for page_num in range(len(pdf_document)):
        #     page = pdf_document.load_page(page_num)
        #     text += page.get_text()
        # return text

    def extract_text_from_excel(self, file_stream: BytesIO) -> str:
        df = pd.read_excel(file_stream)
        return df.to_string(index=False)

    def extract_text_from_txt(self, file_stream: BytesIO) -> str:
        return file_stream.getvalue().decode("utf-8")

    def extract_text_from_powerpoint(self, file_stream: BytesIO) -> str:
        text = ""
        presentation = Presentation(file_stream)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_doc(self, file_stream: BytesIO) -> str:
        text = ""
        doc = docx.Document(file_stream)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
